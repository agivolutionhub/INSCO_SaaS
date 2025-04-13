#!/usr/bin/env python3
from pathlib import Path
import os
import logging
import tempfile
import shutil
import re
import sys
import uuid
import json
import argparse
from typing import List, Optional, Dict, Any

from pptx import Presentation
from fastapi import FastAPI, APIRouter, UploadFile, File, Form, HTTPException, BackgroundTasks
from fastapi.responses import JSONResponse, FileResponse

# Configuración de logging
logger = logging.getLogger("diapos_split")
logger.setLevel(logging.INFO)
if not logger.handlers:
    handler = logging.StreamHandler()
    formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
    handler.setFormatter(formatter)
    logger.addHandler(handler)

# Directorio para almacenamiento
STORAGE_DIR = Path(os.environ.get("STORAGE_DIR", "./storage"))
STORAGE_DIR.mkdir(exist_ok=True, parents=True)

# DEFINIR EL ROUTER COMO VARIABLE GLOBAL igual que en diapos_autofit.py
router = APIRouter(prefix="/api/pptx", tags=["pptx"])

def split_presentation(
    input_file: str, 
    output_dir: Optional[str] = None, 
    slides_per_chunk: int = 20
) -> List[str]:
    input_path = Path(input_file).resolve()
    logger.info(f"Dividiendo presentación: {input_path}")
    
    if not input_path.exists():
        raise FileNotFoundError(f"No se encuentra el archivo: {input_file}")
    
    if input_path.suffix.lower() != '.pptx':
        raise ValueError(f"El archivo debe ser PPTX: {input_file}")
    
    output_dir = Path(output_dir or input_path.parent / f"{input_path.stem}_partes").resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    
    try:
        # Cargar presentación y calcular chunks
        prs = Presentation(input_path)
        total_slides = len(prs.slides)
        num_chunks = (total_slides + slides_per_chunk - 1) // slides_per_chunk
        logger.info(f"Presentación tiene {total_slides} diapositivas, se crearán {num_chunks} archivos")
        
        output_files = []
        base_name = re.sub(r'_(autofit|translated|parte\d*)$', '', input_path.stem)
        
        # Procesar cada chunk
        for chunk in range(num_chunks):
            start_idx = chunk * slides_per_chunk
            end_idx = min((chunk + 1) * slides_per_chunk, total_slides)
            logger.info(f"Procesando parte {chunk+1}: diapositivas {start_idx+1}-{end_idx}")
            
            # Crear copia temporal de la presentación
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                temp_path = temp_file.name
            
            try:
                # Copiar y procesar
                shutil.copy2(input_path, temp_path)
                new_prs = Presentation(temp_path)
                
                # Eliminar diapositivas fuera del rango (orden inverso)
                slides_to_delete = list(range(0, start_idx)) + list(range(end_idx, total_slides))
                for idx in sorted(slides_to_delete, reverse=True):
                    slide_id = new_prs.slides._sldIdLst[idx].rId
                    new_prs.part.drop_rel(slide_id)
                    del new_prs.slides._sldIdLst[idx]
                
                # Guardar el archivo resultante
                output_filename = f"{base_name}_parte {chunk+1}.pptx"
                output_path = output_dir / output_filename
                
                new_prs.save(str(output_path))
                output_files.append(str(output_path))
                
            except Exception as e:
                logger.error(f"Error al procesar parte {chunk+1}: {str(e)}")
                raise
            finally:
                try:
                    os.unlink(temp_path)
                except:
                    pass
                    
        return output_files
        
    except Exception as e:
        logger.error(f"Error durante el proceso: {str(e)}", exc_info=True)
        raise RuntimeError(f"Error al dividir la presentación: {str(e)}")

# ----- PARTE CLI -----
def parse_args():
    parser = argparse.ArgumentParser(description="Divide presentaciones PPTX en archivos más pequeños")
    parser.add_argument("input", type=Path, help="Archivo PPTX a dividir", nargs="?")
    parser.add_argument("-o", "--output-dir", type=Path, help="Directorio para guardar los archivos resultantes")
    parser.add_argument("-s", "--slides", type=int, default=20, help="Número de diapositivas por archivo (default: 20)")
    parser.add_argument("--api", action="store_true", help="Iniciar como servidor API")
    return parser.parse_args()

def cli_main():
    try:
        args = parse_args()
        
        if args.api:
            api_main()
            return 0
            
        if not args.input:
            logger.error("Error: Se requiere archivo de entrada o --api")
            return 1
            
        input_file = args.input
        
        # Validar archivo de entrada
        if not input_file.exists():
            logger.error(f"Error: No se encuentra el archivo {input_file}")
            return 1
        
        if input_file.suffix.lower() != '.pptx':
            logger.error(f"Error: El archivo debe ser PPTX: {input_file}")
            return 1
        
        output_dir = args.output_dir or input_file.parent / f"{input_file.stem}_partes"
        slides_per_chunk = args.slides
        
        logger.info(f"Dividiendo presentación: {input_file}")
        logger.info(f"Diapositivas por archivo: {slides_per_chunk}")
        logger.info(f"Directorio de salida: {output_dir}")
        
        # Dividir la presentación
        output_files = split_presentation(
            input_file=str(input_file),
            output_dir=str(output_dir),
            slides_per_chunk=slides_per_chunk
        )
        
        # Mostrar resultados
        logger.info("\nArchivos generados:")
        for i, file_path in enumerate(output_files, 1):
            file = Path(file_path)
            file_size = file.stat().st_size / (1024 * 1024)  # Tamaño en MB
            logger.info(f"  {i}. {file.name} ({file_size:.1f} MB)")
        
        logger.info(f"\nProceso completado. Se generaron {len(output_files)} archivos en {output_dir}")
        return 0
    
    except Exception as e:
        logger.error(f"Error: {str(e)}")
        return 1

# ----- PARTE API -----
def save_to_storage(file_path: str) -> str:
    file_id = str(uuid.uuid4())
    dest_dir = STORAGE_DIR / file_id
    dest_dir.mkdir(exist_ok=True)
    
    dest_path = dest_dir / Path(file_path).name
    shutil.copy2(file_path, dest_path)
    
    return file_id

def process_pptx_task(input_path: str, output_dir: str, slides_per_chunk: int, job_id: str) -> None:
    result_file = STORAGE_DIR / f"{job_id}_result.json"
    
    try:
        # Comprobar tamaño y procesar archivo
        file_size_mb = os.path.getsize(input_path) / (1024 * 1024)
        logger.info(f"Procesando: {input_path} ({file_size_mb:.2f} MB), diapositivas/chunk: {slides_per_chunk}")
        
        # Dividir la presentación
        output_files = split_presentation(input_path, output_dir, slides_per_chunk)
        
        if not output_files:
            raise Exception("No se generaron archivos de salida")
        
        # Guardar resultados
        results = []
        for output_file in output_files:
            file_id = save_to_storage(str(output_file))
            filename = Path(output_file).name
            results.append({
                "file_id": file_id,
                "filename": filename,
                "url": f"/api/pptx/files/{file_id}/{filename}"
            })
        
        # Guardar resultado exitoso
        with open(result_file, "w", encoding="utf-8") as f:
            json.dump({"status": "completed", "files": results}, f, ensure_ascii=False)
            
    except Exception as e:
        # Guardar resultado con error
        logger.error(f"Error procesando PPTX: {str(e)}", exc_info=True)
        try:
            with open(result_file, "w", encoding="utf-8") as f:
                json.dump({"status": "error", "message": str(e)}, f, ensure_ascii=False)
        except Exception as write_error:
            logger.error(f"Error al escribir resultado: {str(write_error)}")
    finally:
        # Limpiar archivos temporales
        try:
            if os.path.exists(input_path):
                os.unlink(input_path)
            if os.path.exists(output_dir):
                shutil.rmtree(output_dir)
        except Exception as e:
            logger.error(f"Error limpiando temporales: {str(e)}")

# Definir los endpoints directamente en el router global
@router.post("/split")
async def split_pptx_endpoint(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    slides_per_chunk: int = Form(20)
) -> JSONResponse:
    try:
        logger.info(f"Solicitud recibida: archivo={file.filename}, slides_per_chunk={slides_per_chunk}")
        
        # Validar parámetros
        if not file.filename or not file.filename.lower().endswith(".pptx"):
            raise HTTPException(status_code=400, detail="El archivo debe ser PPTX")
        
        if not 1 <= slides_per_chunk <= 100:
            raise HTTPException(status_code=400, detail="El número de diapositivas debe estar entre 1 y 100")
        
        # Crear directorios temporales
        temp_dir = tempfile.mkdtemp()
        input_path = os.path.join(temp_dir, file.filename)
        output_dir = os.path.join(temp_dir, "output")
        os.makedirs(output_dir, exist_ok=True)
        
        # Guardar archivo
        total_size = 0
        with open(input_path, "wb") as buffer:
            chunk_size = 1024 * 1024  # 1MB chunks
            while chunk := await file.read(chunk_size):
                buffer.write(chunk)
                total_size += len(chunk)
        
        logger.info(f"Archivo guardado: {input_path} ({total_size/1024/1024:.2f} MB)")
        
        # Generar ID de trabajo e iniciar tarea
        job_id = str(uuid.uuid4())
        background_tasks.add_task(process_pptx_task, input_path, output_dir, slides_per_chunk, job_id)
        logger.info(f"Tarea iniciada: job_id={job_id}")
        
        # Devolver respuesta inmediata
        return JSONResponse({
            "job_id": job_id,
            "status": "processing",
            "message": "Procesando presentación en segundo plano"
        })
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error inesperado: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Error inesperado: {str(e)}")

@router.get("/jobs/{job_id}")
async def get_job_status(job_id: str) -> JSONResponse:
    result_file = STORAGE_DIR / f"{job_id}_result.json"
    
    if not result_file.exists():
        return JSONResponse({
            "job_id": job_id,
            "status": "processing",
            "message": "El trabajo sigue en proceso"
        })
    
    try:
        with open(result_file, "r", encoding="utf-8") as f:
            return JSONResponse(json.load(f))
    except Exception as e:
        logger.error(f"Error leyendo resultado {result_file}: {str(e)}")
        return JSONResponse({
            "job_id": job_id,
            "status": "error",
            "message": f"Error al recuperar el estado: {str(e)}"
        })

@router.get("/files/{file_id}/{filename}")
async def get_file(file_id: str, filename: str) -> FileResponse:
    file_path = STORAGE_DIR / file_id / filename
    
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="Archivo no encontrado")
    
    return FileResponse(
        path=file_path,
        filename=filename,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

# Función para crear API independiente (mantener para retrocompatibilidad)
def create_api():
    app = FastAPI(title="PPTX Splitter API")
    app.include_router(router)
    return app

# Función para iniciar API independiente
def api_main():
    import uvicorn
    port = int(os.environ.get("PORT", 8000))
    host = os.environ.get("HOST", "0.0.0.0")
    app = create_api()
    logger.info(f"Iniciando API en http://{host}:{port}")
    uvicorn.run(app, host=host, port=port)

# Función para obtener el router, siguiendo el mismo patrón que diapos_autofit.py
get_router = lambda: router

# Punto de entrada
if __name__ == "__main__":
    sys.exit(cli_main())
