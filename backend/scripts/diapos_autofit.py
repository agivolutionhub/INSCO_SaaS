#!/usr/bin/env python3
"""
Script único para aplicar autofit a presentaciones PowerPoint.

Combina la funcionalidad de tres componentes anteriores:
- Servicio de procesamiento (lógica principal)
- CLI para uso desde terminal
- API para uso como servicio web

Ajusta automáticamente el tamaño de texto en presentaciones
para que se adapte a los contenedores.
"""
from pathlib import Path
import sys, os, re, argparse, uuid, shutil, json, time
import logging
from typing import Dict, Any, List, Optional
from fastapi import APIRouter, UploadFile, File, Form, HTTPException, BackgroundTasks
from fastapi.responses import JSONResponse, FileResponse
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE

# Configuración de logging estandarizado
logger = logging.getLogger("autofit")
handler = logging.StreamHandler()
formatter = logging.Formatter('%(levelname)s: %(message)s')
handler.setFormatter(formatter)
logger.addHandler(handler)
logger.setLevel(logging.INFO)

# Configuración por defecto
BASE_DIR = Path(__file__).resolve().parent.parent
DEFAULT_INPUT_DIR = BASE_DIR / "storage/input/diapos"
DEFAULT_OUTPUT_DIR = BASE_DIR / "storage/output/autofit"

# Directorios de almacenamiento
STORAGE_DIR = Path(os.environ.get("AUTOFIT_STORAGE_DIR", BASE_DIR / "storage/autofit")).resolve()

# Asegurar que los directorios existan
for directory in [DEFAULT_INPUT_DIR, DEFAULT_OUTPUT_DIR, STORAGE_DIR]:
    directory.mkdir(parents=True, exist_ok=True)

#===================================
# BLOQUE 1: Funciones de servicio
#===================================

def procesar_pptx(pptx_entrada, pptx_salida=None, silent=False):
    """
    Procesa una presentación PPTX aplicando autofit a todos los textos.
    
    Args:
        pptx_entrada: Ruta al archivo PPTX a procesar
        pptx_salida: Ruta de salida (opcional)
        silent: Si se debe mostrar información por consola
        
    Returns:
        Ruta al archivo procesado
    """
    pptx_entrada = Path(pptx_entrada)
    if not silent:
        logger.info(f"Validando archivo de entrada: {pptx_entrada}")
        
    if not pptx_entrada.exists() or pptx_entrada.suffix.lower() != '.pptx':
        error_msg = f"Archivo inválido o no encontrado: {pptx_entrada}"
        logger.error(error_msg)
        raise FileNotFoundError(error_msg)
    
    pptx_salida = pptx_salida or DEFAULT_OUTPUT_DIR / f"{pptx_entrada.stem}_autofit{pptx_entrada.suffix}"
    pptx_salida = Path(pptx_salida)
    
    if not silent:
        logger.info(f"Creando directorio de salida: {pptx_salida.parent}")
        
    pptx_salida.parent.mkdir(parents=True, exist_ok=True)
    
    try:
        if not silent:
            logger.info(f"Procesando archivo: {pptx_entrada} -> {pptx_salida}")
        
        presentacion = Presentation(pptx_entrada)
        num_slides = len(presentacion.slides)
        
        if not silent:
            logger.info(f"Presentación cargada, {num_slides} diapositivas")
        
        for i, slide in enumerate(presentacion.slides):
            if not silent:
                logger.info(f"Procesando diapositiva {i+1}/{num_slides}")
                
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text_frame"):
                        shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    
                    if hasattr(shape, "has_table") and shape.has_table:
                        for cell in (cell for row in shape.table.rows for cell in row.cells if hasattr(cell, "text_frame")):
                            cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                except Exception as e:
                    if not silent:
                        logger.warning(f"Error al ajustar shape: {e}")
        
        if not silent:
            logger.info(f"Guardando presentación procesada en: {pptx_salida}")
            
        presentacion.save(pptx_salida)
        
        if not silent:
            logger.info(f"✅ Archivo generado: {pptx_salida} ({num_slides} diapositivas)")
        
        # Verificar que el archivo se guardó correctamente
        if not pptx_salida.exists():
            error_msg = f"El archivo no se guardó correctamente: {pptx_salida}"
            logger.error(error_msg)
            raise FileNotFoundError(error_msg)
            
        return str(pptx_salida)
    except Exception as e:
        error_msg = f"Error al procesar archivo: {str(e)}"
        logger.error(error_msg)
        raise Exception(error_msg)

def procesar_lote(directorio=None, salida=None, silent=False):
    """
    Procesa un lote de archivos PPTX aplicando autofit.
    
    Args:
        directorio: Directorio con archivos PPTX a procesar
        salida: Directorio donde guardar los archivos procesados
        silent: Si se debe mostrar información por consola
        
    Returns:
        Diccionario con información de los archivos procesados
    """
    inicio = time.time()
    directorio, salida = Path(directorio or DEFAULT_INPUT_DIR), Path(salida or DEFAULT_OUTPUT_DIR)
    directorio.mkdir(parents=True, exist_ok=True)
    salida.mkdir(parents=True, exist_ok=True)
    
    es_archivo = directorio.is_file() and directorio.suffix.lower() == '.pptx'
    archivos = [directorio] if es_archivo else [f for f in directorio.glob("**/*.pptx") if not f.stem.endswith('_autofit')]
    
    if not archivos:
        if not silent:
            logger.warning(f"No se encontraron archivos PPTX en {directorio}")
        return {"encontrados": 0, "procesados": 0, "errores": 0, "archivos": []}
    
    if not silent:
        logger.info(f"Encontrados {len(archivos)} archivos PPTX para procesar")
    
    resultados = {"encontrados": len(archivos), "procesados": 0, "errores": 0, "archivos": []}
    
    for archivo in archivos:
        try:
            if directorio.is_dir():
                ruta_relativa = archivo.relative_to(directorio)
                destino = salida / ruta_relativa.parent if len(ruta_relativa.parts) > 1 else salida
                destino.mkdir(parents=True, exist_ok=True)
                salida_archivo = destino / f"{archivo.stem}_autofit{archivo.suffix}"
            else:
                salida_archivo = salida / f"{archivo.stem}_autofit{archivo.suffix}"
            
            resultado = procesar_pptx(archivo, salida_archivo, silent=True)
            resultados["procesados"] += 1
            resultados["archivos"].append({
                "entrada": str(archivo),
                "salida": resultado,
                "nombre": archivo.name,
                "estado": "completado"
            })
            
            if not silent:
                logger.info(f"✓ {archivo.name} → {salida_archivo.name}")
        except Exception as e:
            resultados["errores"] += 1
            resultados["archivos"].append({
                "entrada": str(archivo),
                "nombre": archivo.name,
                "estado": "error",
                "mensaje": str(e)
            })
            
            if not silent:
                logger.error(f"✗ Error al procesar {archivo.name}: {str(e)}")
    
    tiempo = time.time() - inicio
    resultados["tiempo"] = tiempo
    
    if not silent:
        logger.info(f"\n📊 Resumen: {resultados['procesados']}/{len(archivos)} archivos procesados ({resultados['errores']} errores) en {tiempo:.2f}s")
        logger.info(f"📁 Archivos generados en: {salida}")
    
    return resultados

#===================================
# BLOQUE 2: Funciones de utilidad
#===================================

def mostrar_menu_carpetas(directorio=DEFAULT_INPUT_DIR):
    """Muestra un menú interactivo para seleccionar carpeta."""
    directorio = Path(directorio)
    directorio.mkdir(parents=True, exist_ok=True)
    subcarpetas = [d for d in directorio.iterdir() if d.is_dir()]
    
    if not subcarpetas:
        logger.warning(f"No se encontraron subcarpetas en {directorio}")
        return None
    
    # Ordenar carpetas por nombre numérico si es posible
    subcarpetas.sort(key=lambda c: int(re.match(r'^(\d+)\.?\s*', c.name).group(1)) 
                    if re.match(r'^(\d+)\.?\s*', c.name) else float('inf'))
    
    print("\nCarpetas disponibles:")
    print("0. [TODAS LAS CARPETAS]")
    for i, carpeta in enumerate(subcarpetas, 1):
        print(f"{i}. {carpeta.name} ({len(list(carpeta.glob('*.pptx')))} archivos PPTX)")
    
    while True:
        try:
            opcion = int(input("\nSelecciona carpeta a procesar [0]: ") or "0")
            if 0 <= opcion <= len(subcarpetas):
                return None if opcion == 0 else subcarpetas[opcion - 1]
            print(f"Opción inválida. Selecciona entre 0 y {len(subcarpetas)}")
        except ValueError:
            print("Ingresa un número válido")

def success_response(data: Any = None, message: Optional[str] = None) -> Dict[str, Any]:
    """Genera una respuesta de éxito estandarizada."""
    response = {"status": "success"}
    if message:
        response["message"] = message
    if data is not None:
        response["data"] = data
    return response

def error_response(message: str) -> Dict[str, Any]:
    """Genera una respuesta de error estandarizada."""
    return {"status": "error", "message": message}

#===================================
# BLOQUE 3: Router FastAPI
#===================================

router = APIRouter(prefix="/api/autofit", tags=["autofit"])

@router.post("/upload-pptx")
async def upload_pptx_for_autofit(file: UploadFile = File(...)):
    """Sube un archivo PPTX para aplicar autofit."""
    try:
        file_id = str(uuid.uuid4())
        filename = file.filename
        original_name = Path(filename).stem
        file_extension = Path(filename).suffix
        
        if file_extension.lower() != '.pptx':
            raise HTTPException(status_code=400, detail=error_response("Solo se permiten archivos PPTX"))
        
        file_location = STORAGE_DIR / f"{file_id}{file_extension}"
        
        with open(file_location, "wb") as f:
            shutil.copyfileobj(file.file, f)
        
        if not file_location.exists():
            raise HTTPException(status_code=500, detail=error_response("No se pudo guardar el archivo"))
        
        return success_response({
            "file_id": file_id,
            "filename": filename,
            "original_name": original_name,
            "file_path": str(file_location)
        })
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=500, 
            detail=error_response(f"Error al procesar la carga del archivo: {str(e)}")
        )

@router.post("/process")
async def process_autofit(file_id: str = Form(...), original_name: str = Form(None)):
    """Procesa un archivo PPTX previamente subido para aplicar autofit."""
    try:
        # Buscar archivo por ID
        files = list(STORAGE_DIR.glob(f"{file_id}.*"))
        
        if not files:
            raise HTTPException(status_code=404, detail=error_response("Archivo no encontrado"))
        
        file_path = files[0]
        
        # Usar el nombre original o un UUID si no está disponible
        output_filename = f"{original_name or file_id}_autofit.pptx"
        output_path = STORAGE_DIR / output_filename
        
        # Procesar archivo
        result_path = procesar_pptx(file_path, output_path, silent=False)
        
        # Verificar que el archivo procesado existe
        if not Path(result_path).exists():
            raise HTTPException(
                status_code=500, 
                detail=error_response("El procesamiento falló: No se generó el archivo")
            )
        
        # Construir respuesta
        response_data = {
            "file_id": file_id,
            "original_name": original_name,
            "processed_file": result_path,
            "output_filename": output_filename,
            "download_url": f"/api/autofit/download/{output_filename}"
        }
        
        return success_response(response_data)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=500, 
            detail=error_response(f"Error al procesar el archivo: {str(e)}")
        )

@router.get("/download/{filename}")
async def download_file(filename: str):
    """Descarga un archivo procesado por autofit."""
    file_path = STORAGE_DIR / filename
    if not file_path.exists():
        raise HTTPException(
            status_code=404, 
            detail=error_response("Archivo no encontrado")
        )
    
    return FileResponse(
        path=file_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=filename
    )

#===================================
# BLOQUE 4: CLI principal
#===================================

def main():
    """Función principal del script."""
    parser = argparse.ArgumentParser(
        description="Aplica autofit a presentaciones PPTX usando python-pptx",
        formatter_class=argparse.ArgumentDefaultsHelpFormatter
    )
    parser.add_argument("pptx_file", nargs="?", type=Path, 
                      help="Archivo PPTX a procesar (opcional)")
    parser.add_argument("-o", "--output", type=Path, 
                      help="Ruta de salida para el archivo procesado")
    parser.add_argument("-i", "--input-dir", type=Path, default=DEFAULT_INPUT_DIR, 
                      help="Directorio de búsqueda de archivos PPTX")
    
    group = parser.add_mutually_exclusive_group()
    group.add_argument("-b", "--batch", action="store_true", 
                     help="Procesar todos los archivos")
    group.add_argument("--interactive", action="store_true", 
                     help="Modo interactivo para seleccionar carpeta")
    
    parser.add_argument("--output-dir", type=Path, default=DEFAULT_OUTPUT_DIR, 
                      help="Directorio para archivos de salida")
    
    args = parser.parse_args()
    
    try:
        if args.interactive or len(sys.argv) <= 1:
            # Modo interactivo
            logger.info("Iniciando modo interactivo")
            carpeta = mostrar_menu_carpetas(args.input_dir)
            
            if carpeta is None:
                if input("¿Procesar TODAS las carpetas? (s/N): ").lower() == 's':
                    logger.info(f"Procesando todos los archivos en {args.input_dir}")
                    procesar_lote(args.input_dir, args.output_dir)
                else:
                    logger.info("Operación cancelada")
            elif input(f"¿Procesar la carpeta {carpeta.name}? (S/n): ").lower() != 'n':
                logger.info(f"Procesando carpeta {carpeta.name}")
                procesar_lote(carpeta, args.output_dir / carpeta.name)
            else:
                logger.info("Operación cancelada")
        elif args.pptx_file or (not args.batch and len(sys.argv) > 1):
            # Procesamiento de archivo único
            archivo = args.pptx_file or Path(sys.argv[1])
            salida = args.output or (Path(sys.argv[2]) if len(sys.argv) > 2 else None)
            logger.info(f"Procesando archivo: {archivo}")
            procesar_pptx(archivo, salida)
        else:
            # Procesamiento por lotes
            logger.info(f"Procesando todos los archivos en {args.input_dir}")
            procesar_lote(args.input_dir, args.output_dir)
            
        return 0
    except Exception as e:
        logger.error(f"Error: {str(e)}")
        return 1

# Para uso como script independiente
if __name__ == "__main__":
    sys.exit(main())

# Para integración con FastAPI, exponer el router
get_autofit_router = lambda: router
