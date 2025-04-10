#!/usr/bin/env python3
from pathlib import Path
import os, logging, tempfile, shutil, re
from pptx import Presentation
from typing import List, Optional

# Configurar logging
logger = logging.getLogger("pptx_service")

def split_presentation(input_file: str, output_dir: Optional[str] = None, slides_per_chunk: int = 20) -> List[str]:
    """
    Divide una presentación PPTX en archivos más pequeños usando python-pptx
    
    Args:
        input_file: Ruta al archivo PPTX a dividir
        output_dir: Directorio donde guardar los archivos resultantes
        slides_per_chunk: Número de diapositivas por archivo
        
    Returns:
        list: Lista de rutas a los archivos PPTX generados
    """
    input_path = Path(input_file).resolve()
    logger.info(f"Dividiendo presentación: {input_path}")
    
    if not input_path.exists():
        raise FileNotFoundError(f"No se encuentra el archivo: {input_file}")
    
    if input_path.suffix.lower() != '.pptx':
        raise ValueError(f"El archivo debe ser PPTX: {input_file}")
    
    output_dir = Path(output_dir or input_path.parent / f"{input_path.stem}_partes").resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    logger.info(f"Archivos de salida se guardarán en: {output_dir}")
    
    try:
        # Cargar presentación y calcular chunks
        prs = Presentation(input_path)
        total_slides = len(prs.slides)
        num_chunks = (total_slides + slides_per_chunk - 1) // slides_per_chunk
        logger.info(f"Presentación tiene {total_slides} diapositivas, se crearán {num_chunks} archivos")
        
        output_files = []
        
        # Limpiar el nombre base quitando sufijos conocidos
        base_name = re.sub(r'_(autofit|translated|parte\d*)$', '', input_path.stem)
        
        # Procesar cada chunk
        for chunk in range(num_chunks):
            start_idx = chunk * slides_per_chunk
            end_idx = min((chunk + 1) * slides_per_chunk, total_slides)
            logger.info(f"Procesando parte {chunk+1}: diapositivas {start_idx+1}-{end_idx}")
            
            # Crear copia temporal de la presentación
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                temp_path = temp_file.name
            
            try:
                # Copiar y procesar
                shutil.copy2(input_path, temp_path)
                new_prs = Presentation(temp_path)
                
                # Eliminar diapositivas fuera del rango (orden inverso)
                slides_to_delete = list(range(0, start_idx)) + list(range(end_idx, total_slides))
                for idx in sorted(slides_to_delete, reverse=True):
                    slide_id = new_prs.slides._sldIdLst[idx].rId
                    new_prs.part.drop_rel(slide_id)
                    del new_prs.slides._sldIdLst[idx]
                
                # Guardar el archivo resultante
                output_filename = f"{base_name}_parte {chunk+1}.pptx"
                output_path = output_dir / output_filename
                
                new_prs.save(str(output_path))
                output_files.append(str(output_path))
                logger.info(f"Guardado: {output_path}")
                
            except Exception as e:
                logger.error(f"Error al procesar parte {chunk+1}: {str(e)}")
                raise
            finally:
                try:
                    os.unlink(temp_path)
                except:
                    pass
        
        # Mostrar resumen de resultados
        if output_files:
            files_info = "\n".join(f"  {i}. {Path(f).name} ({Path(f).stat().st_size/1024/1024:.1f} MB)" 
                                for i, f in enumerate(output_files, 1))
            logger.info(f"Proceso completado. Se generaron {len(output_files)} archivos\n{files_info}")
        
        return output_files
        
    except Exception as e:
        logger.error(f"Error durante el proceso: {str(e)}", exc_info=True)
        raise RuntimeError(f"Error al dividir la presentación: {str(e)}") 