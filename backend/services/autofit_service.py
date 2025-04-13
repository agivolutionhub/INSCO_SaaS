from pathlib import Path
import time
import logging
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE

# Configuración de logging estandarizado
logger = logging.getLogger("autofit-service")
if not logger.handlers:
    handler = logging.StreamHandler()
    formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
    handler.setFormatter(formatter)
    logger.addHandler(handler)
logger.setLevel(logging.INFO)

# Configuración por defecto
BASE_DIR = Path(__file__).resolve().parent.parent.parent
DEFAULT_INPUT_DIR = BASE_DIR / "data/input/diapos"
DEFAULT_OUTPUT_DIR = BASE_DIR / "data/output/01_autofit2"

# Asegurar que los directorios por defecto existan
DEFAULT_INPUT_DIR.mkdir(parents=True, exist_ok=True)
DEFAULT_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

def procesar_pptx(pptx_entrada, pptx_salida=None, silent=False):
    """
    Procesa una presentación PPTX aplicando autofit a todos los textos.
    
    Args:
        pptx_entrada: Ruta al archivo PPTX a procesar
        pptx_salida: Ruta de salida (opcional)
        silent: Si se debe mostrar información por consola
        
    Returns:
        Ruta al archivo procesado
        
    Raises:
        FileNotFoundError: Si el archivo de entrada no existe
        Exception: Para cualquier otro error durante el procesamiento
    """
    pptx_entrada = Path(pptx_entrada)
    if not silent:
        logger.info(f"Validando archivo de entrada: {pptx_entrada}")
        
    if not pptx_entrada.exists() or pptx_entrada.suffix.lower() != '.pptx':
        error_msg = f"Archivo inválido o no encontrado: {pptx_entrada}"
        logger.error(error_msg)
        raise FileNotFoundError(error_msg)
    
    pptx_salida = pptx_salida or DEFAULT_OUTPUT_DIR / f"{pptx_entrada.stem}_autofit{pptx_entrada.suffix}"
    pptx_salida = Path(pptx_salida)
    
    if not silent:
        logger.info(f"Creando directorio de salida: {pptx_salida.parent}")
        
    pptx_salida.parent.mkdir(parents=True, exist_ok=True)
    
    try:
        if not silent:
            logger.info(f"Procesando archivo: {pptx_entrada} -> {pptx_salida}")
        
        presentacion = Presentation(pptx_entrada)
        num_slides = len(presentacion.slides)
        
        if not silent:
            logger.info(f"Presentación cargada, {num_slides} diapositivas")
        
        for i, slide in enumerate(presentacion.slides):
            if not silent:
                logger.info(f"Procesando diapositiva {i+1}/{num_slides}")
                
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text_frame"):
                        shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    
                    if hasattr(shape, "has_table") and shape.has_table:
                        for cell in (cell for row in shape.table.rows for cell in row.cells if hasattr(cell, "text_frame")):
                            cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                except Exception as e:
                    if not silent:
                        logger.warning(f"Error al ajustar shape: {e}")
        
        if not silent:
            logger.info(f"Guardando presentación procesada en: {pptx_salida}")
            
        presentacion.save(pptx_salida)
        
        if not silent:
            logger.info(f"✅ Archivo generado: {pptx_salida} ({num_slides} diapositivas)")
        
        # Verificar que el archivo se guardó correctamente
        if not pptx_salida.exists():
            error_msg = f"El archivo no se guardó correctamente: {pptx_salida}"
            logger.error(error_msg)
            raise FileNotFoundError(error_msg)
            
        return str(pptx_salida)
    except Exception as e:
        error_msg = f"Error al procesar archivo: {str(e)}"
        logger.error(error_msg, exc_info=not silent)
        raise Exception(error_msg)

def procesar_lote(directorio=None, salida=None, silent=False):
    """
    Procesa un lote de archivos PPTX aplicando autofit.
    
    Args:
        directorio: Directorio con archivos PPTX a procesar
        salida: Directorio donde guardar los archivos procesados
        silent: Si se debe mostrar información por consola
        
    Returns:
        Diccionario con información de los archivos procesados
    """
    inicio = time.time()
    directorio, salida = Path(directorio or DEFAULT_INPUT_DIR), Path(salida or DEFAULT_OUTPUT_DIR)
    directorio.mkdir(parents=True, exist_ok=True)
    salida.mkdir(parents=True, exist_ok=True)
    
    es_archivo = directorio.is_file() and directorio.suffix.lower() == '.pptx'
    archivos = [directorio] if es_archivo else [f for f in directorio.glob("**/*.pptx") if not f.stem.endswith('_autofit')]
    
    if not archivos:
        if not silent:
            logger.warning(f"No se encontraron archivos PPTX en {directorio}")
        return {"encontrados": 0, "procesados": 0, "errores": 0, "archivos": []}
    
    if not silent:
        logger.info(f"Encontrados {len(archivos)} archivos PPTX para procesar")
    
    resultados = {"encontrados": len(archivos), "procesados": 0, "errores": 0, "archivos": []}
    
    for archivo in archivos:
        try:
            if directorio.is_dir():
                ruta_relativa = archivo.relative_to(directorio)
                destino = salida / ruta_relativa.parent if len(ruta_relativa.parts) > 1 else salida
                destino.mkdir(parents=True, exist_ok=True)
                salida_archivo = destino / f"{archivo.stem}_autofit{archivo.suffix}"
            else:
                salida_archivo = salida / f"{archivo.stem}_autofit{archivo.suffix}"
            
            resultado = procesar_pptx(archivo, salida_archivo, silent=True)
            resultados["procesados"] += 1
            resultados["archivos"].append({
                "entrada": str(archivo),
                "salida": resultado,
                "nombre": archivo.name,
                "estado": "completado"
            })
            
            if not silent:
                logger.info(f"✓ {archivo.name} → {salida_archivo.name}")
        except Exception as e:
            resultados["errores"] += 1
            resultados["archivos"].append({
                "entrada": str(archivo),
                "nombre": archivo.name,
                "estado": "error",
                "mensaje": str(e)
            })
            
            if not silent:
                logger.error(f"✗ Error al procesar {archivo.name}: {str(e)}")
    
    tiempo = time.time() - inicio
    resultados["tiempo"] = tiempo
    
    if not silent:
        logger.info(f"\n📊 Resumen: {resultados['procesados']}/{len(archivos)} archivos procesados ({resultados['errores']} errores) en {tiempo:.2f}s")
        logger.info(f"📁 Archivos generados en: {salida}")
    
    return resultados 